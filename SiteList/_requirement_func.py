#-*- coding:utf-8 -*-

from bs4 import BeautifulSoup
from requests import get

from SiteList._translate_j2k import *

import codecs

from re import *

from click import clear

from time import sleep

from json import loads

# from concurrent.futures import ThreadPoolExecutor
from os import remove

from PIL import Image

import pytesseract

from telepot.exception import TelegramError

from urllib3.exceptions import ProtocolError

from urllib.parse import urlparse

from datetime import datetime

from json import dumps

from os import chdir, system, rename

from pprint import pprint


dirLoc = '.\\jp_trans_tmp\\'

pytesseract.pytesseract.tesseract_cmd = r'C:\\Users\\power\\AppData\\Local\\Tesseract-OCR\\tesseract.exe'


banner = '''
   ___        _    _      _           _   _                _ 
  |_  |      | |  | |    | |         | \ | |              | |
    | |______| |  | | ___| |__ ______|  \| | _____   _____| |
    | |______| |/\| |/ _ \ '_ \______| . ` |/ _ \ \ / / _ \ |
/\__/ /      \  /\  /  __/ |_) |     | |\  | (_) \ V /  __/ |
\____/        \/  \/ \___|_.__/      \_| \_/\___/ \_/ \___|_|                                                           
        _____                   _       _                           
       |_   _|                 | |     | |                          
         | |_ __ __ _ _ __  ___| | __ _| |_ ___  _ __               
         | | '__/ _` | '_ \/ __| |/ _` | __/ _ \| '__|              
         | | | | (_| | | | \__ \ | (_| | || (_) | |                 
         \_/_|  \__,_|_| |_|___/_|\__,_|\__\___/|_|

                * 일본 웹 소설 번역기 *  
'''


def GetSoup(url, referer):
    headers = {
        'referer': referer,
        "User-Agent": "Mozilla/5.0",
    }
    while True:
        try:
            html = get(url, headers=headers, cookies={'over18':'yes'}).text
            soup = BeautifulSoup(html, 'html.parser')
            return soup
        except:
            sleep(2)



def WriteFile(text: str, filename: str):
    f = codecs.open(filename, mode='w', encoding='utf-8')
    f.write(u'{}'.format(text))
    f.close()



def PrintProgressbar(iteration, total, prefix='', suffix='', decimals=1, length=10, fill='#'):
    percent = ("{0:." + str(decimals) + "f}").format(100 * (iteration / float(total)))
    filledLength = int(length * iteration // total)
    bar = fill * filledLength + ' ' * (length - filledLength)
    print(f'\r{prefix} |{bar}| {percent}% {suffix}', end='\r')
    if iteration == total: 
        print()


def GetFileName(filename: str):
    toReplace = {
        '\\':'', '/':'', ':':'-', '\"':'', '-':'_',
        '?':'', '<':'[', '>':']', '|':'-', '*':'',
        '\n':'', '\t':'', ':':'-', ' ':' '
    }

    for key, value in toReplace.items():
        filename = str(filename).replace(key, value)

    return filename


def ReplacingText(text:str, repl_dict: dict):
    for key, value in repl_dict.items():
        replaced_text = str(text).replace(key, value)

    return replaced_text


def PrettyJson(msg):
    return dumps(msg, indent=4, sort_keys=True, ensure_ascii=False)


def MessageLog(msg):
    now = datetime.now()
    log = "[{0}/{1}/{2} - {3}:{4}:{5}]\n{6}\n\n\n".format(now.year, now.month, now.day, now.hour, now.minute, now.second, str(PrettyJson(msg)))

    f = codecs.open('./message_log.log', 'a+', encoding='utf-8')
    f.write(log)
    f.close()



def SendLongMessage(msg:str, chat_id, bot):
    j_text = t_j2k(msg)
    
    length = 4000
    cut = [j_text[i:i+length] for i in range(0, len(j_text), length)]

    for c in cut:
        bot.sendMessage(chat_id, c)



def GetTextInAudio(audio_msg, chat_id, bot):
    from speech_recognition import AudioFile, Recognizer, UnknownValueError

    rec = Recognizer()

    file_id = audio_msg['file_id']

    tempAudio = dirLoc+file_id+'.'+audio_msg['mime_type'].split('/')[1]

    bot.download_file(file_id, tempAudio)

    from pydub import AudioSegment
    filename = file_id + '.wav'
    file_loc = dirLoc + filename

    sound = AudioSegment.from_file(tempAudio)
    sound.export(file_loc, format='wav')


    with AudioFile(file_loc) as AudioSrc:
        content = rec.record(AudioSrc)
    try:    
        text = rec.recognize_google(
            audio_data=content,
            language='ja-JP',
        )
    except UnknownValueError:
        bot.sendMessage(chat_id, 'None')
        return

    mainText = t_j2k(text)
    bot.sendMessage(chat_id, mainText)

    remove(tempAudio)
    remove(file_loc)




def GetTextInDocument(document_msg, chat_id, bot):
    mime_type = document_msg['mime_type']

    file_id = document_msg['file_id']
    filename = document_msg['file_name']

    file_loc = dirLoc + filename

    bot.download_file(file_id, file_loc)

    t_title = t_j2k(filename)


    if 'image/' in mime_type:
        img = Image.open(file_loc)
        j_text = pytesseract.image_to_string(img, lang='jpn')
        print(j_text)
        text = sub(r'[\s+]', '', j_text)
        bot.sendMessage(chat_id, t_j2k(text))


    elif mime_type == 'text/plain':
        f = codecs.open(file_loc, 'r', encoding='utf-8').read()
        
        t_content = ''
        for t in f.split('\n'):
            t_content += t_j2k(t)

        WriteFile(t_content, dirLoc+t_title)
        bot.sendDocument(chat_id, codecs.open(dirLoc+t_title, 'r', encoding='utf-8'))

        remove(dirLoc+t_title)


    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        from docx import Document

        document = Document(file_loc)

        for p in document.paragraphs:
            inline = p.runs

            for i in inline:
                print(i.text)
                i.text = t_j2k(i.text)

        document.save(dirLoc+t_title)
        bot.sendDocument(chat_id, codecs.open(dirLoc+t_title, 'rb'))
        
        remove(dirLoc+t_title)
        

    
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        from pptx import Presentation
        
        prs = Presentation(file_loc)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame

                    for para in text_frame.paragraphs:
                        for r in para.runs:
                            r.text = t_j2k(r.text)

        prs.save(dirLoc+t_title)
        bot.sendDocument(chat_id, codecs.open(dirLoc+t_title, 'rb'))
        remove(dirLoc+t_title)
        


    # elif mime_type == '':
    #     pass


    else:
        bot.sendMessage(chat_id, '아직 다른 파일은 지원하지 않아요!')
        return





def GetUrlInLogfile(logFile:str):
    f = codecs.open(logFile, 'r')
    r = f.readlines()
    f.close()
    return r






def GetTextInWebpage(url, chat_id, bot):

    o = urlparse(url)

    import random
    import string

    chars = "".join([random.choice(string.ascii_lowercase) for i in range(15)])

    logFile = f'{dirLoc}{chars}.wget_log'
    system(f'.\\Utils\\wget.exe -o {logFile} -P {dirLoc} --convert-links --html-extension -p -k {url}')

    l = GetUrlInLogfile(logFile)

    for i, s in enumerate(l):
        if '[text/html]' in s:
            html_loc = l[i+1].split('Saving to: ')[1].replace('\n', '').replace("'", '')
            break

    remove(logFile)

    f = codecs.open(html_loc, 'r', encoding='utf-8')
    html = f.read()
    f.close()

    remove(html_loc)

    s = str(BeautifulSoup(html, 'html.parser').prettify()).split('\n')

    isScript = False
    isStyle = False


    newHtml = ''

    for h in s:

        rp_h = h.replace(' ', '')

        if rp_h == '':
            newHtml += '\n'

        elif rp_h.startswith('<script'):
            print('스크립트 시작')
            isScript = True
            newHtml += h


        elif rp_h.startswith('</script>'):
            print("스크립트 종료")
            isScript = False
            newHtml += h


        elif rp_h.startswith('<style'):
            print("스타일 시작")
            isStyle = True
            newHtml += h


        elif rp_h.startswith('</style>'):
            print('스타일 종료')
            isStyle = False
            newHtml += h


        elif rp_h.startswith('<img' or '<embed' or '<source'):
            print("이미지, 동영상, 파일")
            newHtml += h


        else:
            if isScript or isStyle:
                newHtml += h

            else:
                if ((rp_h.startswith('"') or rp_h == sub('<.*?>', '', rp_h))):
                    if bool(findall('[\u2E80-\u9FFF]', rp_h)):
                        newHtml += t_j2k(h, isHtml=True)
                        print('텍스트 번역: ', h)
                else:
                    newHtml += h


    newHtml = str(BeautifulSoup(newHtml, 'html.parser').prettify())


    if html_loc[-5:] != '.html':
        html_loc += '.html'

    WriteFile(newHtml, html_loc)

    hostedURL = 'https://kimdrew.loca.lt/' + html_loc.replace('./', '')

    print(hostedURL)

    bot.sendMessage(chat_id, hostedURL)
    











# def t_executor(func, *args, **kwargs):
#     with ThreadPoolExecutor(max_workers=64) as executor:
#         future = executor.submit(func, *args, **kwargs)
    
#         return future.result()


# def TagRemover(html):
#     r = compile('<.*?>|&([a-z0-9]+|#[0-9]{1,6}|#x[0-9a-f]{1,6});')
#     return sub(r, '', html)